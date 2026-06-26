from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "AI_Fake_Job_Offer_Detector_Normal_Presentation.pptx"
ASSETS = ROOT / "presentation_assets"
ASSETS.mkdir(exist_ok=True)

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "bg": RGBColor(238, 242, 244),
    "panel": RGBColor(255, 255, 255),
    "ink": RGBColor(23, 32, 33),
    "muted": RGBColor(96, 111, 115),
    "line": RGBColor(210, 221, 224),
    "accent": RGBColor(15, 118, 110),
    "accent2": RGBColor(36, 91, 155),
    "red": RGBColor(184, 50, 58),
    "amber": RGBColor(196, 122, 17),
    "dark": RGBColor(17, 33, 35),
}


def font(size=18, bold=False, color=None):
    return {"size": Pt(size), "bold": bold, "color": color or COLORS["ink"]}


def set_run(run, style):
    run.font.name = "Aptos"
    run.font.size = style["size"]
    run.font.bold = style["bold"]
    run.font.color.rgb = style["color"]


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    for run in p.runs:
        set_run(run, font(size, bold, color))
    return box


def add_multiline(slide, lines, x, y, w, h, size=17, color=None, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if bullet:
            p.text = f"- {line}"
        for run in p.runs:
            set_run(run, font(size, False, color or COLORS["ink"]))
        p.space_after = Pt(7)
    return box


def add_panel(slide, x, y, w, h, fill=None, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or COLORS["panel"]
    shape.line.color.rgb = line or COLORS["line"]
    shape.line.width = Pt(1)
    return shape


def add_badge(slide, text, x, y, w, color):
    badge = add_panel(slide, x, y, w, 0.38, fill=color, line=color)
    add_text(slide, text, x, y + 0.06, w, 0.22, size=10, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    return badge


def add_arrow(slide, x1, y1, x2, y2, color=None):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color or COLORS["accent"]
    line.line.width = Pt(2.4)
    line.line.end_arrowhead = True
    return line


def background(slide, title, subtitle=None, section="JobShield"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["accent"]
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(1.05), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = COLORS["dark"]
    side.line.fill.background()
    add_text(slide, "JD", 0.25, 0.28, 0.55, 0.34, size=14, bold=True, color=RGBColor(216, 251, 239), align=PP_ALIGN.CENTER)
    add_text(slide, section, 0.18, 6.95, 0.72, 0.22, size=8, bold=True, color=RGBColor(185, 201, 203), align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.38, 0.45, 8.4, 0.55, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1.4, 1.02, 9.8, 0.28, size=11, color=COLORS["muted"])


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["dark"]
    add_text(slide, "AI-Based Fake Job Offer Detector", 0.85, 0.92, 8.7, 0.8, size=36, bold=True, color=RGBColor(255, 255, 255))
    add_text(slide, "Using AI methods, cryptography, network security, and cyber law awareness", 0.9, 1.78, 8.7, 0.35, size=16, color=RGBColor(202, 219, 221))
    add_panel(slide, 0.9, 2.62, 4.9, 2.8, fill=RGBColor(23, 49, 51), line=RGBColor(56, 89, 91))
    add_multiline(slide, ["Fake offer detection", "Domain and link inspection", "SHA-256 document fingerprint", "Cybercrime reporting guidance"], 1.25, 3.0, 4.2, 1.7, size=18, color=RGBColor(237, 247, 246), bullet=True)
    add_panel(slide, 6.35, 1.45, 5.85, 4.8, fill=RGBColor(238, 242, 244), line=RGBColor(62, 96, 98))
    slide.shapes.add_picture(str(make_ui_image()), Inches(6.62), Inches(1.75), width=Inches(5.3))
    add_badge(slide, "PROJECT PRESENTATION", 0.9, 6.42, 2.25, COLORS["accent"])


def make_ui_image():
    path = ASSETS / "ui_mock.png"
    img = Image.new("RGB", (1050, 650), (238, 242, 244))
    d = ImageDraw.Draw(img)
    try:
        title = ImageFont.truetype("arialbd.ttf", 28)
        normal = ImageFont.truetype("arial.ttf", 20)
        small = ImageFont.truetype("arial.ttf", 16)
    except OSError:
        title = normal = small = None
    d.rounded_rectangle((30, 30, 255, 620), radius=18, fill=(17, 33, 35))
    d.text((62, 68), "JobShield", fill=(237, 247, 246), font=title)
    for i, item in enumerate(["Home", "Fake Job Detection", "History", "About"]):
        y = 145 + i * 62
        fill = (38, 74, 73) if i == 1 else (17, 33, 35)
        d.rounded_rectangle((55, y, 230, y + 42), radius=8, fill=fill, outline=(60, 92, 94))
        d.text((72, y + 11), item, fill=(235, 245, 245), font=small)
    d.rounded_rectangle((290, 55, 1000, 600), radius=18, fill=(255, 255, 255), outline=(210, 221, 224), width=2)
    d.text((330, 92), "Fake Job Detection", fill=(23, 32, 33), font=title)
    d.rounded_rectangle((330, 150, 650, 500), radius=12, fill=(248, 250, 251), outline=(210, 221, 224))
    d.text((360, 180), "Offer Evidence", fill=(23, 32, 33), font=normal)
    d.rounded_rectangle((360, 230, 620, 266), radius=6, fill=(255, 255, 255), outline=(210, 221, 224))
    d.rounded_rectangle((360, 292, 620, 420), radius=6, fill=(255, 255, 255), outline=(210, 221, 224))
    d.rounded_rectangle((690, 150, 960, 500), radius=12, fill=(248, 250, 251), outline=(210, 221, 224))
    d.text((720, 182), "Risk score", fill=(96, 111, 115), font=small)
    d.text((720, 218), "87%", fill=(184, 50, 58), font=title)
    d.rounded_rectangle((720, 290, 920, 316), radius=13, fill=(255, 226, 224))
    d.text((740, 296), "Likely Fake", fill=(146, 51, 51), font=small)
    img.save(path)
    return path


def content_slide(prs, title, subtitle, bullets, visual=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, title, subtitle)
    add_panel(slide, 1.35, 1.65, 5.45, 4.95)
    add_multiline(slide, bullets, 1.72, 2.02, 4.75, 3.95, size=18, bullet=True)
    if visual:
        visual(slide)
    return slide


def visual_problem(slide):
    slide.shapes.add_picture(str(make_email_image()), Inches(7.15), Inches(1.72), width=Inches(5.55))


def make_email_image():
    path = ASSETS / "fake_email_mock.png"
    img = Image.new("RGB", (980, 620), (246, 249, 250))
    d = ImageDraw.Draw(img)
    try:
        title = ImageFont.truetype("arialbd.ttf", 30)
        normal = ImageFont.truetype("arial.ttf", 22)
        small = ImageFont.truetype("arial.ttf", 18)
    except OSError:
        title = normal = small = None
    d.rounded_rectangle((40, 45, 940, 575), radius=18, fill=(255, 255, 255), outline=(210, 221, 224), width=3)
    d.rectangle((40, 45, 940, 118), fill=(17, 33, 35))
    d.text((72, 68), "Inbox - Job Offer", fill=(237, 247, 246), font=title)
    d.text((76, 146), "From: hr.companyjobs@gmail.com", fill=(184, 50, 58), font=normal)
    d.text((76, 185), "Subject: Direct Joining - No Interview Required", fill=(23, 32, 33), font=normal)
    y = 245
    lines = [
        "Congratulations! You are selected for Software Engineer.",
        "Pay refundable security deposit within 24 hours.",
        "Send Aadhaar, PAN, bank statement on WhatsApp.",
        "Payment link: http://bit.ly/job-fee-payment",
    ]
    for line in lines:
        d.text((96, y), line, fill=(55, 65, 68), font=small)
        y += 48
    d.rounded_rectangle((74, 470, 350, 525), radius=12, fill=(255, 226, 224), outline=(238, 180, 176))
    d.text((98, 486), "High risk indicators", fill=(146, 51, 51), font=normal)
    img.save(path)
    return path


def visual_modules(slide):
    labels = [("AI Text Analysis", 7.35, 1.85, COLORS["accent"]), ("Network Security", 10.0, 1.85, COLORS["accent2"]), ("SHA-256 Hash", 7.35, 4.05, COLORS["amber"]), ("Cyber Law Guide", 10.0, 4.05, COLORS["red"])]
    for text, x, y, c in labels:
        add_panel(slide, x, y, 2.15, 1.05, fill=RGBColor(255, 255, 255), line=c)
        add_text(slide, text, x + 0.12, y + 0.36, 1.9, 0.25, size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_panel(slide, 8.78, 3.0, 2.2, 0.82, fill=COLORS["dark"], line=COLORS["dark"])
    add_text(slide, "Risk Score", 9.08, 3.27, 1.6, 0.25, size=15, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    for x1, y1, x2, y2 in [(8.42, 2.9, 9.15, 3.1), (10.05, 2.9, 10.15, 3.1), (8.42, 4.05, 9.15, 3.8), (10.05, 4.05, 10.15, 3.8)]:
        add_arrow(slide, x1, y1, x2, y2)


def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, "System Workflow", "How the detector converts offer evidence into a final decision")
    steps = [
        ("User input", "Text or .txt file"),
        ("Extract indicators", "Emails, URLs, risk terms"),
        ("Run checks", "AI rules + network logic"),
        ("Hash document", "SHA-256 fingerprint"),
        ("Final output", "Score, verdict, guidance"),
    ]
    x = 1.45
    for i, (head, sub) in enumerate(steps):
        add_panel(slide, x, 2.65, 1.9, 1.25, fill=COLORS["panel"], line=COLORS["accent"] if i % 2 == 0 else COLORS["accent2"])
        add_text(slide, head, x + 0.13, 2.98, 1.62, 0.24, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x + 0.13, 3.3, 1.62, 0.24, size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.93, 3.28, x + 2.38, 3.28)
        x += 2.35
    add_badge(slide, "End-to-end detection flow", 4.88, 5.25, 2.7, COLORS["accent"])


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, "Architecture Diagram", "Client-side prototype with modular detection logic")
    add_panel(slide, 1.55, 2.75, 2.2, 1.0, fill=RGBColor(255, 255, 255), line=COLORS["accent"])
    add_text(slide, "Frontend UI", 1.9, 3.1, 1.45, 0.25, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, 5.3, 1.55, 2.4, 1.0, fill=RGBColor(255, 255, 255), line=COLORS["accent2"])
    add_text(slide, "Text Analyzer", 5.63, 1.9, 1.7, 0.25, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, 5.3, 3.05, 2.4, 1.0, fill=RGBColor(255, 255, 255), line=COLORS["accent2"])
    add_text(slide, "Network Checks", 5.55, 3.4, 1.9, 0.25, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, 5.3, 4.55, 2.4, 1.0, fill=RGBColor(255, 255, 255), line=COLORS["accent2"])
    add_text(slide, "Crypto Hash", 5.72, 4.9, 1.55, 0.25, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, 9.3, 2.75, 2.25, 1.0, fill=COLORS["dark"], line=COLORS["dark"])
    add_text(slide, "Risk Engine", 9.65, 3.1, 1.55, 0.25, size=15, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_panel(slide, 9.3, 4.55, 2.25, 1.0, fill=RGBColor(255, 255, 255), line=COLORS["red"])
    add_text(slide, "History Store", 9.63, 4.9, 1.6, 0.25, size=15, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)
    for y in [2.05, 3.55, 5.05]:
        add_arrow(slide, 3.78, 3.25, 5.25, y)
        add_arrow(slide, 7.75, y, 9.28, 3.25)
    add_arrow(slide, 10.42, 3.78, 10.42, 4.52, COLORS["red"])


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, "Conclusion", "Summary of the project's impact and achievements")
    add_panel(slide, 1.35, 1.65, 10.65, 4.95)
    add_multiline(slide, [
        "The AI-Based Fake Job Offer Detector successfully combines multiple cybersecurity approaches to protect job seekers.",
        "Key achievements include rule-based AI text analysis, network security checks, SHA-256 document hashing, and cyber law guidance.",
        "The prototype demonstrates how client-side technologies can provide immediate risk assessment without requiring backend infrastructure.",
        "Future enhancements could include machine learning models, live API integrations, and verified employer portals.",
        "This project contributes to digital literacy and cyber safety awareness among job seekers in India and globally."
    ], 1.72, 2.02, 9.95, 3.95, size=18, bullet=True)
    add_badge(slide, "Thank You", 5.5, 6.0, 2.0, COLORS["accent"])


def build():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    title_slide(prs)
    content_slide(prs, "Problem Statement", "Understanding the threat of fake job offers", [
        "Fake job offers are a growing cyber threat targeting job seekers worldwide",
        "Scammers use sophisticated impersonation techniques to appear legitimate",
        "Victims often lose money through fraudulent fees, deposits, or data theft",
        "Fresh graduates and unemployed individuals are particularly vulnerable",
        "Lack of awareness and easy verification tools exacerbates the problem",
        "Economic impact includes financial losses and damaged trust in job markets",
        "Cybercrime statistics show increasing incidents of employment fraud"
    ], visual_problem)
    content_slide(prs, "Objectives", "Goals and scope of the project", [
        "Develop an accessible tool for detecting suspicious job offers",
        "Implement multiple detection methods: AI text analysis, network checks, cryptography",
        "Provide clear risk scoring with explainable reasons",
        "Include cyber law guidance for appropriate response actions",
        "Create a browser-based prototype for easy deployment and use",
        "Support both text input and file upload for flexibility",
        "Maintain user privacy by processing data client-side only"
    ])
    content_slide(prs, "Technologies Used", "Technical stack and implementation details", [
        "Frontend: HTML5, CSS3, JavaScript ES6+ for responsive web interface",
        "Web APIs: Web Crypto API for SHA-256 hashing, LocalStorage for data persistence",
        "Text Processing: Regular expressions and rule-based pattern matching",
        "Security: URL parsing, domain validation, link analysis",
        "UI Framework: Custom CSS with light/dark theme support",
        "Build Tools: Python with python-pptx for presentation generation",
        "Image Generation: PIL/Pillow for mock UI screenshots",
        "Deployment: Static web server (Python http.server) for local hosting"
    ], visual_modules)
    architecture_slide(prs)
    flow_slide(prs)
    content_slide(prs, "SHA-256 Cryptography Method", "Document integrity verification using cryptographic hashing", [
        "SHA-256 is a cryptographic hash function producing 256-bit (32-byte) hash values",
        "Used to create unique fingerprints of job offer documents",
        "Any change in input text produces completely different hash output",
        "Provides document integrity verification without revealing content",
        "Implemented using browser's Web Crypto API for client-side processing",
        "Hash values are displayed in hexadecimal format for readability",
        "Helps users verify if offer content has been tampered with",
        "Cryptographically secure and collision-resistant algorithm"
    ])
    content_slide(prs, "Advantages and Limitations", "Strengths and areas for improvement", [
        "Advantages:",
        "Browser-based, no installation required",
        "Client-side processing ensures privacy",
        "Multiple detection methods for comprehensive analysis",
        "Explainable AI approach with clear reasoning",
        "Cyber law guidance for appropriate actions",
        "Lightweight and fast performance",
        "Limitations:",
        "Rule-based detection may miss sophisticated scams",
        "No machine learning model for advanced pattern recognition",
        "Limited to text analysis, no image or PDF processing",
        "Cannot verify actual company hiring practices",
        "Depends on user-provided evidence quality"
    ])
    content_slide(prs, "Future Scope", "Potential enhancements and research directions", [
        "Integrate machine learning models trained on real fake job datasets",
        "Add support for PDF and image OCR processing",
        "Implement live API checks for domain reputation and WHOIS data",
        "Create verified employer portal with digital signatures",
        "Develop mobile app version for broader accessibility",
        "Add multi-language support for global users",
        "Implement real-time URL scanning and malware detection",
        "Partner with employment agencies for official verification services",
        "Research advanced NLP techniques for better text analysis",
        "Explore blockchain-based verification systems"
    ])
    conclusion_slide(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()