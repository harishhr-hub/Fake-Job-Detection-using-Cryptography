from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "AI_Fake_Job_Offer_Detector_Detailed_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(243, 247, 248)
INK = RGBColor(25, 35, 38)
MUTED = RGBColor(91, 107, 112)
LINE = RGBColor(211, 222, 226)
PANEL = RGBColor(255, 255, 255)
DARK = RGBColor(17, 33, 35)
TEAL = RGBColor(15, 118, 110)
BLUE = RGBColor(38, 91, 155)
RED = RGBColor(184, 50, 58)
AMBER = RGBColor(196, 122, 17)
GREEN = RGBColor(33, 132, 88)
PALE_TEAL = RGBColor(224, 246, 241)
PALE_RED = RGBColor(255, 231, 231)
PALE_BLUE = RGBColor(229, 240, 255)
PALE_AMBER = RGBColor(255, 244, 220)


def rgb(hex_value):
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_text(paragraph, size=16, bold=False, color=INK, align=None):
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def textbox(slide, text, x, y, w, h, size=16, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    set_text(p, size=size, bold=bold, color=color, align=align)
    return box


def bullets(slide, items, x, y, w, h, size=15, color=INK, numbered=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        prefix = f"{index + 1}. " if numbered else "- "
        p.text = f"{prefix}{item}"
        p.space_after = Pt(7)
        set_text(p, size=size, color=color)
    return box


def panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1)
    return s


def chip(slide, text, x, y, w, fill, color=RGBColor(255, 255, 255)):
    panel(slide, x, y, w, 0.38, fill=fill, line=fill)
    textbox(slide, text, x + 0.04, y + 0.09, w - 0.08, 0.18, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.2):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    ln.line.end_arrowhead = True
    return ln


def slide_base(prs, title, subtitle=None, tag="JobShield"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.9), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = DARK
    side.line.fill.background()

    textbox(slide, "JD", 0.19, 0.38, 0.5, 0.25, size=12, bold=True, color=RGBColor(216, 251, 239), align=PP_ALIGN.CENTER)
    textbox(slide, tag, 0.12, 7.04, 0.65, 0.18, size=7, bold=True, color=RGBColor(188, 205, 208), align=PP_ALIGN.CENTER)
    textbox(slide, title, 1.25, 0.42, 9.6, 0.45, size=24, bold=True)
    if subtitle:
        textbox(slide, subtitle, 1.27, 0.93, 10.7, 0.25, size=11, color=MUTED)
    return slide


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
    panel(slide, 0.85, 0.7, 11.65, 6.15, fill=rgb("173033"), line=rgb("34585b"))
    textbox(slide, "AI-Based Fake Job Offer Detector", 1.25, 1.08, 8.4, 0.58, size=33, bold=True, color=RGBColor(255, 255, 255))
    textbox(slide, "Using AI methods, cryptography, network security, and cyber law awareness", 1.28, 1.76, 8.3, 0.32, size=15, color=RGBColor(209, 226, 228))
    chip(slide, "PROJECT PRESENTATION", 1.28, 2.32, 2.28, TEAL)

    cards = [
        ("AI Text Analysis", "Detects fraud language and suspicious hiring promises.", 1.35, 3.12, PALE_TEAL, TEAL),
        ("Network Security", "Checks emails, domains, HTTP links, and short URLs.", 4.08, 3.12, PALE_BLUE, BLUE),
        ("SHA-256 Hashing", "Creates an integrity fingerprint for submitted offer text.", 6.81, 3.12, PALE_AMBER, AMBER),
        ("Cyber Law Guide", "Guides evidence preservation and fraud reporting steps.", 9.54, 3.12, PALE_RED, RED),
    ]
    for title, desc, x, y, fill, color in cards:
        panel(slide, x, y, 2.33, 1.92, fill=fill, line=color)
        textbox(slide, title, x + 0.16, y + 0.22, 1.98, 0.25, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x + 0.17, y + 0.72, 1.98, 0.72, size=10.5, color=INK, align=PP_ALIGN.CENTER)
    textbox(slide, "Complete project explanation with clear workflows, methods, diagrams, and future scope", 1.3, 6.0, 9.5, 0.35, size=14, color=RGBColor(216, 231, 233))


def problem_slide(prs):
    slide = slide_base(prs, "Problem Statement", "Fake job offers use impersonation, payment traps, and phishing links to target job seekers.")
    panel(slide, 1.25, 1.45, 4.25, 4.9)
    textbox(slide, "Why fake job offers are dangerous", 1.55, 1.78, 3.3, 0.28, size=17, bold=True)
    bullets(slide, [
        "Fraudsters copy real company names, logos, and HR titles.",
        "Candidates are asked to pay deposits, training fees, or verification fees.",
        "Personal documents such as Aadhaar, PAN, and bank details can be misused.",
        "Shortened links and fake domains can redirect users to phishing or payment pages.",
        "Freshers and urgent job seekers are especially vulnerable.",
    ], 1.55, 2.25, 3.6, 2.7, size=14)

    textbox(slide, "Typical Scam Path", 7.35, 1.42, 3.2, 0.3, size=17, bold=True)
    steps = [("Fake HR message", RED), ("Attractive salary", AMBER), ("Urgent payment", RED), ("Data theft / loss", DARK)]
    x = 6.1
    for i, (name, color) in enumerate(steps):
        panel(slide, x, 2.55, 1.6, 1.15, fill=PANEL, line=color)
        textbox(slide, name, x + 0.12, 2.96, 1.34, 0.22, size=12.5, bold=True, color=color, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(slide, x + 1.62, 3.12, x + 2.0, 3.12, RED)
        x += 2.0
    chip(slide, "Goal: detect before damage happens", 7.35, 5.28, 3.1, TEAL)


def objectives_slide(prs):
    slide = slide_base(prs, "Project Objectives", "The system converts raw offer evidence into an understandable security decision.")
    objs = [
        ("Detect", "Identify suspicious language, payment requests, and unrealistic promises.", TEAL),
        ("Verify", "Check recruiter email domains, URLs, and document integrity hash.", BLUE),
        ("Explain", "Show clear reasons so users understand why an offer is risky.", AMBER),
        ("Guide", "Provide safe next steps and cybercrime reporting awareness.", RED),
    ]
    x = 1.35
    for title, desc, color in objs:
        panel(slide, x, 1.65, 2.55, 3.15, fill=PANEL, line=color)
        textbox(slide, title, x + 0.25, 2.0, 2.05, 0.35, size=21, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x + 0.25, 2.65, 2.05, 1.1, size=14, color=MUTED, align=PP_ALIGN.CENTER)
        x += 2.85
    panel(slide, 2.0, 5.35, 9.8, 0.75, fill=PALE_TEAL, line=TEAL)
    textbox(slide, "Final objective: help job seekers make safer decisions before replying, paying money, or sharing documents.", 2.25, 5.58, 9.3, 0.24, size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def tech_slide(prs):
    slide = slide_base(prs, "Technologies Used", "Each technology supports a specific part of the detection pipeline.")
    rows = [
        ("HTML, CSS, JavaScript", "Builds the web app, screens, forms, and result panels.", TEAL),
        ("Rule-based AI/NLP logic", "Scores suspicious text patterns such as fees, urgency, and no-interview claims.", BLUE),
        ("Regex + URL parsing", "Extracts recruiter emails, domains, and links from offer text.", AMBER),
        ("Web Crypto API", "Generates SHA-256 hash for document integrity.", TEAL),
        ("LocalStorage", "Stores detection history and theme preference in the browser.", BLUE),
        ("Cyber law references", "Guides reporting and evidence-preservation awareness.", RED),
    ]
    y = 1.45
    for tech, use, color in rows:
        panel(slide, 1.35, y, 3.3, 0.62, fill=PANEL, line=color)
        textbox(slide, tech, 1.58, y + 0.18, 2.8, 0.18, size=12.5, bold=True, color=color)
        panel(slide, 4.85, y, 6.95, 0.62, fill=PANEL, line=LINE)
        textbox(slide, use, 5.08, y + 0.17, 6.45, 0.2, size=12.5, color=INK)
        y += 0.78


def architecture_slide(prs):
    slide = slide_base(prs, "System Architecture", "Modular client-side architecture used by the current prototype.")
    layers = [
        ("User Interface", "Home, Detection, History, About, Theme Toggle", 1.3, TEAL),
        ("Input Processing", "Offer text, company domain, file upload, extracted emails/URLs", 2.5, BLUE),
        ("Detection Engine", "Text pattern scoring + network checks + SHA-256 hashing", 3.7, AMBER),
        ("Output Layer", "Risk score, verdict, findings, hash, cyber law guidance", 4.9, RED),
    ]
    for name, desc, y, color in layers:
        panel(slide, 2.0, y, 9.5, 0.78, fill=PANEL, line=color)
        textbox(slide, name, 2.28, y + 0.23, 2.15, 0.2, size=14, bold=True, color=color)
        textbox(slide, desc, 4.75, y + 0.22, 6.35, 0.22, size=13, color=INK)
    for y1 in [2.08, 3.28, 4.48]:
        arrow(slide, 6.75, y1, 6.75, y1 + 0.4, TEAL)


def workflow_slide(prs):
    slide = slide_base(prs, "Complete Workflow", "Step-by-step path from user input to final decision.")
    steps = [
        ("1", "Enter offer details", "Company name, domain, text, or .txt file"),
        ("2", "Extract indicators", "Emails, URLs, risk words, missing details"),
        ("3", "Run security checks", "Free-mail, domain mismatch, short links, HTTP links"),
        ("4", "Calculate hash", "SHA-256 fingerprint of submitted text"),
        ("5", "Generate result", "Score, verdict, findings, law guidance"),
        ("6", "Save history", "Local record for later review"),
    ]
    x_positions = [1.25, 3.25, 5.25, 7.25, 9.25, 11.25]
    for i, (num, title, desc) in enumerate(steps):
        panel(slide, x_positions[i], 2.2, 1.55, 2.1, fill=PANEL, line=TEAL if i % 2 == 0 else BLUE)
        textbox(slide, num, x_positions[i] + 0.55, 2.42, 0.45, 0.32, size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        textbox(slide, title, x_positions[i] + 0.12, 2.98, 1.3, 0.34, size=11.5, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x_positions[i] + 0.13, 3.42, 1.28, 0.56, size=8.7, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(slide, x_positions[i] + 1.57, 3.25, x_positions[i + 1] - 0.05, 3.25)
    chip(slide, "Clear sequential flow for viva explanation", 4.85, 5.55, 3.65, TEAL)


def input_slide(prs):
    slide = slide_base(prs, "Input Evidence Collected", "The quality of detection depends on the evidence provided by the user.")
    panel(slide, 1.35, 1.55, 4.95, 4.6)
    textbox(slide, "User provides", 1.65, 1.87, 2.3, 0.28, size=18, bold=True)
    bullets(slide, [
        "Company name mentioned in offer.",
        "Official company domain, for example company.com.",
        "Offer text copied from email, chat, or document.",
        "Optional plain-text file upload.",
    ], 1.65, 2.35, 4.1, 1.8, size=15)
    panel(slide, 7.0, 1.55, 4.3, 4.6, fill=rgb("fbfcfc"), line=LINE)
    textbox(slide, "Example suspicious input", 7.35, 1.88, 2.6, 0.28, size=18, bold=True, color=RED)
    bullets(slide, [
        "From: hr.companyjobs@gmail.com",
        "Direct joining without interview",
        "Pay refundable security deposit",
        "Send PAN, Aadhaar, bank statement",
        "Payment link: http://bit.ly/job-fee",
    ], 7.35, 2.38, 3.6, 2.1, size=14, color=INK)


def text_analysis_slide(prs):
    slide = slide_base(prs, "AI Text Analysis Method", "Rule-based AI-style scoring makes the prototype explainable.")
    categories = [
        ("Payment Demand", "security deposit, registration fee, training fee", RED),
        ("Urgency Pressure", "within 24 hours, limited seats, act fast", AMBER),
        ("Unrealistic Hiring", "no interview, direct joining, guaranteed job", RED),
        ("Sensitive Data", "Aadhaar, PAN, bank statement, OTP, password", BLUE),
        ("Unofficial Contact", "WhatsApp only, Telegram, personal email", TEAL),
    ]
    y = 1.55
    for category, examples, color in categories:
        panel(slide, 1.35, y, 3.1, 0.58, fill=PANEL, line=color)
        textbox(slide, category, 1.58, y + 0.17, 2.5, 0.18, size=12.5, bold=True, color=color)
        panel(slide, 4.62, y, 6.95, 0.58, fill=PANEL, line=LINE)
        textbox(slide, examples, 4.85, y + 0.17, 6.35, 0.18, size=12.5, color=INK)
        y += 0.78
    textbox(slide, "Why this method is useful: every detected signal is shown as a readable finding, so the result is not a black box.", 1.55, 6.1, 10.0, 0.28, size=14, bold=True, color=TEAL)


def network_slide(prs):
    slide = slide_base(prs, "Network Security Checks", "The system reviews email and URL signals commonly seen in phishing scams.")
    checks = [
        ("Free email domain", "hr.company@gmail.com", "Risk: recruiter is not using official company mail"),
        ("Domain mismatch", "official: company.com, email: company-careers.net", "Risk: impersonation domain"),
        ("HTTP link", "http://fake-job-form.com", "Risk: insecure link"),
        ("Shortened URL", "bit.ly/job-fee-payment", "Risk: destination hidden"),
        ("Payment link", "/pay, /upi, /checkout", "Risk: fee collection trap"),
    ]
    y = 1.42
    for title, example, reason in checks:
        panel(slide, 1.25, y, 2.5, 0.7, fill=PALE_BLUE, line=BLUE)
        textbox(slide, title, 1.45, y + 0.23, 2.1, 0.18, size=12.5, bold=True, color=BLUE)
        panel(slide, 3.95, y, 3.5, 0.7, fill=PANEL, line=LINE)
        textbox(slide, example, 4.15, y + 0.23, 3.1, 0.18, size=11.5, color=INK)
        panel(slide, 7.65, y, 4.15, 0.7, fill=PALE_RED, line=RED)
        textbox(slide, reason, 7.85, y + 0.2, 3.75, 0.24, size=11.2, color=RED)
        y += 0.88


def crypto_slide(prs):
    slide = slide_base(prs, "Cryptography Method: SHA-256 Hashing", "Hashing gives the submitted offer a document-integrity fingerprint.")
    panel(slide, 1.35, 1.55, 4.55, 4.55)
    textbox(slide, "Original offer text", 1.7, 1.9, 2.4, 0.25, size=16, bold=True)
    textbox(slide, "Joining Date: 20 May 2026\nCTC: 6 LPA\nDomain: company.com", 1.7, 2.45, 3.6, 0.9, size=15, color=INK)
    arrow(slide, 5.95, 3.15, 7.05, 3.15, TEAL)
    panel(slide, 7.15, 1.55, 4.55, 4.55)
    textbox(slide, "SHA-256 hash output", 7.5, 1.9, 2.6, 0.25, size=16, bold=True, color=TEAL)
    textbox(slide, "9f86d081884c7d659a2feaa0c55ad015...", 7.5, 2.55, 3.6, 0.55, size=16, color=BLUE)
    textbox(slide, "If one character changes, the hash changes. This helps prove whether the offer content remained unchanged.", 7.5, 3.55, 3.45, 0.8, size=14, color=MUTED)
    chip(slide, "Integrity, not identity verification", 7.55, 5.25, 3.0, AMBER)


def scoring_slide(prs):
    slide = slide_base(prs, "Risk Scoring Method", "Every suspicious indicator adds risk points to the final score.")
    items = [
        ("Security deposit request", 22, RED),
        ("Payment words or UPI/bank transfer", 18, RED),
        ("No interview / direct joining", 18, RED),
        ("Free-mail recruiter domain", 16, AMBER),
        ("Urgent pressure language", 12, AMBER),
        ("Shortened URL", 12, AMBER),
        ("Missing formal offer details", 8, BLUE),
    ]
    y = 1.45
    for label, score, color in items:
        textbox(slide, label, 1.35, y + 0.08, 3.1, 0.2, size=12.5, bold=True)
        panel(slide, 4.6, y, 5.4, 0.36, fill=rgb("e5ecef"), line=rgb("e5ecef"), radius=False)
        panel(slide, 4.6, y, 5.4 * score / 24, 0.36, fill=color, line=color, radius=False)
        textbox(slide, f"+{score}", 10.25, y + 0.04, 0.65, 0.2, size=12.5, bold=True, color=color)
        y += 0.62
    panel(slide, 10.95, 2.55, 1.45, 1.45, fill=RED, line=RED)
    textbox(slide, "87%", 11.15, 2.93, 1.0, 0.32, size=25, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    textbox(slide, "Likely Fake", 11.05, 3.35, 1.25, 0.18, size=10.5, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def output_slide(prs):
    slide = slide_base(prs, "Output and Result Explanation", "The result screen is designed to be understandable for non-technical users.")
    panel(slide, 1.35, 1.45, 3.0, 4.9, fill=PANEL, line=RED)
    textbox(slide, "Risk Score", 1.7, 1.85, 1.4, 0.25, size=15, bold=True, color=MUTED)
    textbox(slide, "87%", 1.7, 2.35, 1.4, 0.55, size=34, bold=True, color=RED)
    chip(slide, "Likely Fake", 1.7, 3.1, 1.5, RED)
    panel(slide, 4.75, 1.45, 3.2, 4.9)
    textbox(slide, "Key Findings", 5.05, 1.85, 1.8, 0.25, size=16, bold=True)
    bullets(slide, ["Payment requested before joining", "Free-mail recruiter domain", "Urgency pressure", "Shortened payment link"], 5.05, 2.35, 2.45, 1.8, size=13)
    panel(slide, 8.35, 1.45, 3.4, 4.9)
    textbox(slide, "Security Evidence", 8.65, 1.85, 1.9, 0.25, size=16, bold=True)
    bullets(slide, ["Network indicators", "SHA-256 document hash", "Cyber law next steps", "History record"], 8.65, 2.35, 2.45, 1.8, size=13)


def history_slide(prs):
    slide = slide_base(prs, "Detection History", "History helps users compare previous scans and preserve a basic local record.")
    panel(slide, 1.25, 1.55, 10.8, 4.65)
    headers = ["Company", "Domain", "Score", "Verdict", "Saved Evidence"]
    xs = [1.55, 3.55, 5.7, 7.0, 8.85]
    widths = [1.6, 1.7, 0.8, 1.4, 2.7]
    for x, w, h in zip(xs, widths, headers):
        panel(slide, x, 1.95, w, 0.5, fill=DARK, line=DARK, radius=False)
        textbox(slide, h, x + 0.04, 2.1, w - 0.08, 0.14, size=9.5, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    rows = [
        ["ABC Tech", "abc-careers.net", "82%", "Likely Fake", "Preview + SHA-256 hash"],
        ["RealSoft", "realsoft.com", "28%", "Lower Risk", "Preview + SHA-256 hash"],
        ["QuickJobs", "gmail.com", "74%", "Likely Fake", "Preview + SHA-256 hash"],
    ]
    y = 2.55
    for row in rows:
        for x, w, cell in zip(xs, widths, row):
            panel(slide, x, y, w, 0.55, fill=rgb("f8fbfc"), line=LINE, radius=False)
            color = RED if "Likely" in cell or cell in ["82%", "74%"] else INK
            textbox(slide, cell, x + 0.04, y + 0.18, w - 0.08, 0.14, size=9.2, color=color, align=PP_ALIGN.CENTER)
        y += 0.72
    textbox(slide, "Stored in browser LocalStorage. It is simple for demo use and does not require a backend database.", 1.55, 5.55, 9.8, 0.24, size=13.5, bold=True, color=TEAL)


def ui_slide(prs):
    slide = slide_base(prs, "User Interface Screens", "The UI is organized into clear navigation options requested for the project.")
    screens = [
        ("Home", "Explains project purpose and key capabilities.", TEAL),
        ("Fake Job Detection", "Main form for entering offer evidence and viewing results.", BLUE),
        ("Detection History", "Shows recent scans saved in the browser.", AMBER),
        ("About", "Explains objective, law context, and project disclaimer.", RED),
        ("Theme Switch", "Light and dark mode for better usability.", DARK),
    ]
    y = 1.48
    for title, desc, color in screens:
        panel(slide, 1.45, y, 2.75, 0.75, fill=PANEL, line=color)
        textbox(slide, title, 1.7, y + 0.23, 2.2, 0.18, size=13, bold=True, color=color)
        panel(slide, 4.55, y, 6.55, 0.75, fill=PANEL, line=LINE)
        textbox(slide, desc, 4.8, y + 0.22, 6.05, 0.2, size=12.5, color=INK)
        y += 0.9


def law_slide(prs):
    slide = slide_base(prs, "Cyber Law and Reporting Guidance", "The app includes practical guidance for users who find a suspicious offer.")
    panel(slide, 1.35, 1.55, 4.7, 4.75, fill=PALE_RED, line=RED)
    textbox(slide, "Possible cyber offences", 1.7, 1.9, 2.6, 0.28, size=17, bold=True, color=RED)
    bullets(slide, ["Online cheating", "Identity theft", "Phishing", "Impersonation", "Financial fraud"], 1.75, 2.42, 3.0, 1.8, size=15, color=INK)
    panel(slide, 6.55, 1.55, 4.95, 4.75, fill=PANEL, line=TEAL)
    textbox(slide, "Recommended response", 6.9, 1.9, 2.6, 0.28, size=17, bold=True, color=TEAL)
    bullets(slide, [
        "Do not pay any fee or deposit.",
        "Preserve screenshots, emails, headers, phone numbers, URLs, and payment IDs.",
        "Report through cybercrime.gov.in or local cyber cell.",
        "Inform the real company if its brand is misused.",
    ], 6.9, 2.42, 4.0, 2.1, size=14)


def demo_slide(prs):
    slide = slide_base(prs, "Demo Walkthrough", "Use this sequence while presenting the working application.")
    steps = [
        ("Open Home", "Introduce the purpose of JobShield."),
        ("Go to Detection", "Enter company name, official domain, and offer text."),
        ("Upload sample", "Use samples/fake_offer.txt for a quick fake-offer demo."),
        ("Analyze offer", "Explain score, verdict, findings, network checks, and hash."),
        ("Open History", "Show that the detection result was stored locally."),
    ]
    y = 1.45
    for i, (title, desc) in enumerate(steps):
        panel(slide, 1.45, y, 1.0, 0.6, fill=TEAL, line=TEAL)
        textbox(slide, str(i + 1), 1.78, y + 0.2, 0.32, 0.16, size=13, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        panel(slide, 2.75, y, 8.3, 0.6, fill=PANEL, line=LINE)
        textbox(slide, title, 3.02, y + 0.18, 1.7, 0.17, size=12.5, bold=True, color=TEAL)
        textbox(slide, desc, 4.8, y + 0.18, 5.8, 0.17, size=12.2, color=INK)
        y += 0.82


def benefits_limitations_slide(prs):
    slide = slide_base(prs, "Advantages and Limitations", "A clear explanation of what the prototype does well and what can be improved.")
    panel(slide, 1.35, 1.55, 4.95, 4.9, fill=PALE_TEAL, line=TEAL)
    textbox(slide, "Advantages", 1.72, 1.9, 2.0, 0.3, size=18, bold=True, color=TEAL)
    bullets(slide, [
        "Easy to use and runs in the browser.",
        "Fast risk scoring with readable reasons.",
        "No backend required for the current version.",
        "Combines AI-style analysis, security checks, hashing, and law awareness.",
        "Useful for academic demo and basic candidate awareness.",
    ], 1.72, 2.4, 3.85, 2.3, size=14)
    panel(slide, 6.85, 1.55, 4.95, 4.9, fill=PALE_AMBER, line=AMBER)
    textbox(slide, "Limitations", 7.22, 1.9, 2.0, 0.3, size=18, bold=True, color=AMBER)
    bullets(slide, [
        "Rule-based detection is not a trained ML model yet.",
        "No live WHOIS or domain-age lookup in this prototype.",
        "No PDF/image OCR support yet.",
        "Cannot confirm actual company hiring records.",
        "Hashing proves integrity, not whether a company issued the offer.",
    ], 7.22, 2.4, 3.85, 2.3, size=14)


def future_slide(prs):
    slide = slide_base(prs, "Future Scope", "Planned improvements can make the system stronger and closer to production.")
    roadmap = [
        ("Phase 1", "PDF + OCR", "Read offer letters and screenshots directly.", TEAL),
        ("Phase 2", "ML Model", "Train classifier using real and fake offer datasets.", BLUE),
        ("Phase 3", "Live APIs", "WHOIS, SSL, URL reputation, email-header checks.", AMBER),
        ("Phase 4", "Verified Portal", "Companies sign offer letters with QR/digital proof.", RED),
    ]
    x = 1.35
    for phase, title, desc, color in roadmap:
        panel(slide, x, 2.0, 2.45, 2.8, fill=PANEL, line=color)
        chip(slide, phase, x + 0.35, 2.32, 1.28, color)
        textbox(slide, title, x + 0.28, 3.02, 1.85, 0.26, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x + 0.25, 3.55, 1.9, 0.65, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        if x < 9:
            arrow(slide, x + 2.48, 3.38, x + 2.85, 3.38, color)
        x += 2.85


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
    textbox(slide, "Conclusion", 1.1, 0.82, 5.0, 0.55, size=33, bold=True, color=RGBColor(255, 255, 255))
    textbox(slide, "The project demonstrates how multiple cybersecurity ideas can work together to protect job seekers from fake employment offers.", 1.12, 1.58, 9.7, 0.55, size=18, color=RGBColor(215, 230, 232))
    cols = [("Detect", "Suspicious text, links, and emails", TEAL), ("Verify", "Domain match and SHA-256 integrity", BLUE), ("Explain", "Risk score with clear findings", AMBER), ("Report", "Cyber law and evidence guidance", RED)]
    x = 1.15
    for title, desc, color in cols:
        panel(slide, x, 3.05, 2.55, 1.75, fill=rgb("173033"), line=color)
        textbox(slide, title, x + 0.22, 3.43, 2.1, 0.3, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(slide, desc, x + 0.25, 3.98, 2.05, 0.38, size=11.5, color=RGBColor(220, 234, 236), align=PP_ALIGN.CENTER)
        x += 2.85
    textbox(slide, "Thank You", 1.15, 6.22, 2.8, 0.42, size=24, bold=True, color=RGBColor(255, 255, 255))


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)
    problem_slide(prs)
    objectives_slide(prs)
    tech_slide(prs)
    architecture_slide(prs)
    workflow_slide(prs)
    text_analysis_slide(prs)
    network_slide(prs)
    crypto_slide(prs)
    scoring_slide(prs)
    output_slide(prs)
    history_slide(prs)
    ui_slide(prs)
    law_slide(prs)
    demo_slide(prs)
    benefits_limitations_slide(prs)
    future_slide(prs)
    conclusion_slide(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
